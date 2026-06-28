import os
from docx import Document
from pptx import Presentation

def extract_docx(file_path):
    doc = Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())
    return "\n".join(text)

def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides):
        text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text.strip())
    return "\n".join(text)

if __name__ == "__main__":
    base_dir = "/Users/avipriyaghosh/Desktop/telecom-anomaly-system"
    
    print("EXTRACTING DOCX...")
    docx_path = os.path.join(base_dir, "MCA Dissertation- template.docx")
    try:
        print(extract_docx(docx_path))
    except Exception as e:
        print("Error reading docx:", e)
        
    print("\n\nEXTRACTING PPTX 2...")
    pptx2_path = os.path.join(base_dir, "Telecom_Revenue_Leakage_Review_2.pptx")
    try:
        print(extract_pptx(pptx2_path))
    except Exception as e:
        print("Error reading pptx 2:", e)
